#!/usr/bin/env python3
"""
Write a real PowerPoint deck for a course, from its guide.

    python3 pipeline/slides.py econ
    python3 pipeline/slides.py --all

One deck per course, one unit per section: a unit title slide, then each card as
a question slide followed by its answer. That order is the point — a deck that
shows the answer with the question is a document, not a study aid.

Output: app/public/decks/<course>.pptx, which PowerPoint, Keynote and Google
Slides all open. The app links to it from the guide's Slides mode; the in-app
deck is the same content drawn live.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from guide_reader import ROOT, all_courses, read_guide, unit_title  # noqa: E402

# The app's own palette, so a printed deck and the screen are the same object.
INK = RGBColor(0x0A, 0x0B, 0x0E)
PAPER = RGBColor(0xE8, 0xEA, 0xEE)
DIM = RGBColor(0x94, 0x9B, 0xA7)
ACCENT = RGBColor(0xC8, 0xCE, 0xD8)

WIDE = Inches(13.333)
TALL = Inches(7.5)


def blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = INK
    return slide


def text(slide, body: str, *, top: float, size: int, color=PAPER, bold=False, left=0.9, width=11.5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.0))
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = body
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def rule(slide, top: float, width: float = 11.5, color=DIM):
    line = slide.shapes.add_shape(1, Inches(0.9), Inches(top), Inches(width), Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False


def deck(course: str) -> Path:
    guide = read_guide(course)
    prs = Presentation()
    prs.slide_width = WIDE
    prs.slide_height = TALL

    # ── cover ────────────────────────────────────────────────────────────
    slide = blank(prs)
    text(slide, guide["code"].upper(), top=2.4, size=54, bold=True)
    text(slide, guide["name"], top=3.5, size=24, color=DIM)
    text(slide, guide["blurb"], top=4.2, size=16, color=DIM)
    rule(slide, 5.3)
    text(
        slide,
        f"{len(guide['units'])} units · "
        f"{sum(len(u['cards']) for u in guide['units'])} cards · from {guide['source']}",
        top=5.5,
        size=12,
        color=DIM,
    )

    for i, unit in enumerate(guide["units"]):
        title = unit_title(unit["name"])

        slide = blank(prs)
        text(slide, f"UNIT {i + 1}", top=2.6, size=14, color=ACCENT, bold=True)
        text(slide, title, top=3.1, size=40, bold=True)
        text(slide, f"{len(unit['cards'])} things to know", top=4.5, size=16, color=DIM)
        rule(slide, 4.4)

        for n, card in enumerate(unit["cards"], 1):
            q = blank(prs)
            text(q, f"{title.upper()} · {n} OF {len(unit['cards'])}", top=0.6, size=11, color=ACCENT)
            text(q, card["q"], top=2.4, size=32, bold=True)
            text(q, "Answer it before you advance.", top=6.3, size=12, color=DIM)

            a = blank(prs)
            text(a, card["q"], top=0.9, size=18, color=DIM)
            rule(a, 1.9)
            text(a, card["a"], top=2.2, size=22)

    if guide["terms"]:
        slide = blank(prs)
        text(slide, "TERMS", top=0.6, size=14, color=ACCENT, bold=True)
        top = 1.4
        for term in guide["terms"][:10]:
            text(slide, term["t"], top=top, size=16, bold=True)
            text(slide, term["d"], top=top + 0.35, size=13, color=DIM)
            top += 1.0
            if top > 6.4:
                break

    if guide["selfTest"]:
        slide = blank(prs)
        text(slide, "SELF-TEST", top=2.6, size=14, color=ACCENT, bold=True)
        text(slide, "Answer these out loud", top=3.1, size=36, bold=True)
        text(slide, f"{len(guide['selfTest'])} questions, answers on the slide after each", top=4.4, size=15, color=DIM)
        for n, card in enumerate(guide["selfTest"], 1):
            q = blank(prs)
            text(q, f"SELF-TEST · {n} OF {len(guide['selfTest'])}", top=0.6, size=11, color=ACCENT)
            text(q, card["q"], top=2.6, size=30, bold=True)
            a = blank(prs)
            text(a, card["q"], top=0.9, size=17, color=DIM)
            rule(a, 1.9)
            text(a, card["a"], top=2.2, size=21)

    out = ROOT / "app/public/decks" / f"{course}.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"{out}  {len(prs.slides)} slides")
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("course", nargs="?")
    ap.add_argument("--all", action="store_true", help="Every course in the catalog.")
    args = ap.parse_args()

    courses = all_courses() if args.all or not args.course else [args.course]
    for course in courses:
        deck(course)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
